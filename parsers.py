"""
Текстовые парсеры для входных материалов (преподаватель) и работ студентов.

Поддерживаемые форматы (без OCR):
- PDF (через pypdf или PyPDF2 - что установлено; нет ничего - исключение)
- DOCX (python-docx)
- TXT/MD/CSV/TSV/JSON/PY и прочие "текстовые" расширения
- IPYNB (извлекает markdown + code cells)
- ZIP (распаковывает и извлекает текст из поддерживаемых файлов внутри)
"""

from __future__ import annotations

import csv
import hashlib
import io
import json
import mimetypes
import os
import re
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, Optional

# Optional dependencies


_PDF_BACKEND = None
try:
    from pypdf import PdfReader  # type: ignore

    _PDF_BACKEND = "pypdf"
except Exception:
    try:
        from PyPDF2 import PdfReader  # type: ignore

        _PDF_BACKEND = "PyPDF2"
    except Exception:
        PdfReader = None  # type: ignore

_DOCX_BACKEND = None
try:
    from docx import Document as DocxDocument  # type: ignore

    _DOCX_BACKEND = "python-docx"
except Exception:
    DocxDocument = None  # type: ignore

_PPTX_BACKEND = None
try:
    from pptx import Presentation  # type: ignore

    _PPTX_BACKEND = "python-pptx"
except Exception:
    Presentation = None  # type: ignore

# Public constants


OWNER_STUDENT = "student"
OWNER_TEACHER = "teacher"
OWNER_SYSTEM = "system"
OWNER_UNKNOWN = "unknown"

STAGE_TOPIC_ALIGNMENT = "topic_alignment"
STAGE_METHODICS = "methodics"
STAGE_SUBMISSION = "submission"
STAGE_GENERAL = "general"

DEFAULT_MAX_TEXT_CHARS_PER_DOC = 30000
DEFAULT_MAX_TOTAL_CONTEXT_CHARS = 120000
DEFAULT_MAX_ARCHIVE_ENTRIES = 80
DEFAULT_MAX_ARCHIVE_DEPTH = 2
DEFAULT_MAX_FILE_SIZE_BYTES = 25 * 1024 * 1024  # 25 MB


# Data models


@dataclass
class ParsedDocument:
    filename: str
    extension: str
    text: str
    owner: str = OWNER_UNKNOWN
    stage: str = STAGE_GENERAL
    parser_name: str = "unknown"
    mime_type: Optional[str] = None
    source_name: Optional[str] = None
    archive_path: Optional[str] = None
    size_bytes: int = 0
    sha256: str = ""
    warnings: list[str] = field(default_factory=list)
    is_supported: bool = True

    @property
    def doc_id(self) -> str:
        digest = self.sha256[:12] if self.sha256 else "nohash"
        return f"{self.filename}:{digest}"

    @property
    def display_name(self) -> str:
        if self.archive_path:
            return f"{self.archive_path}!{self.filename}"
        return self.filename

    def short_meta(self) -> str:
        parts = [
            f"файл: {self.display_name}",
            f"владелец: {self.owner}",
            f"этап: {self.stage}",
            f"тип: {self.extension or 'unknown'}",
            f"парсер: {self.parser_name}",
        ]
        if self.warnings:
            parts.append("предупреждения: " + "; ".join(self.warnings))
        return " | ".join(parts)

    def to_context_block(self, max_chars: Optional[int] = None) -> str:
        text = self.text or ""
        if max_chars is not None and len(text) > max_chars:
            text = text[:max_chars].rstrip() + "\n...[обрезано]"

        header_lines = [
            f"[ИСТОЧНИК]",
            f"Файл: {self.display_name}",
            f"Владелец: {self.owner}",
            f"Этап: {self.stage}",
            f"Парсер: {self.parser_name}",
        ]
        if self.warnings:
            header_lines.append("Предупреждения: " + "; ".join(self.warnings))

        return "\n".join(header_lines) + "\n\n" + (text or "[пустой текст]")

    def to_dict(self) -> dict[str, Any]:
        return {
            "doc_id": self.doc_id,
            "filename": self.filename,
            "display_name": self.display_name,
            "extension": self.extension,
            "owner": self.owner,
            "stage": self.stage,
            "parser_name": self.parser_name,
            "mime_type": self.mime_type,
            "source_name": self.source_name,
            "archive_path": self.archive_path,
            "size_bytes": self.size_bytes,
            "sha256": self.sha256,
            "warnings": self.warnings,
            "is_supported": self.is_supported,
            "text": self.text,
        }


# Public API


def parse_uploaded_file(
        uploaded_file: Any,
        *,
        owner: str = OWNER_UNKNOWN,
        stage: str = STAGE_GENERAL,
        max_archive_entries: int = DEFAULT_MAX_ARCHIVE_ENTRIES,
        max_archive_depth: int = DEFAULT_MAX_ARCHIVE_DEPTH,
        max_file_size_bytes: int = DEFAULT_MAX_FILE_SIZE_BYTES,
) -> list[ParsedDocument]:
    """
    Универсальный вход для Streamlit UploadedFile / bytes / path / file-like объекта.
    Возвращает список документов, так как zip-архив может распаковаться в несколько текстовых файлов.
    """
    raw_bytes, filename, mime_type, source_name = _coerce_input(uploaded_file)
    return parse_bytes(
        raw_bytes,
        filename=filename,
        mime_type=mime_type,
        owner=owner,
        stage=stage,
        source_name=source_name,
        max_archive_entries=max_archive_entries,
        max_archive_depth=max_archive_depth,
        max_file_size_bytes=max_file_size_bytes,
    )


def parse_uploaded_files(
        uploaded_files: Iterable[Any],
        *,
        owner: str = OWNER_UNKNOWN,
        stage: str = STAGE_GENERAL,
        max_archive_entries: int = DEFAULT_MAX_ARCHIVE_ENTRIES,
        max_archive_depth: int = DEFAULT_MAX_ARCHIVE_DEPTH,
        max_file_size_bytes: int = DEFAULT_MAX_FILE_SIZE_BYTES,
) -> list[ParsedDocument]:
    docs: list[ParsedDocument] = []
    for item in uploaded_files:
        docs.extend(
            parse_uploaded_file(
                item,
                owner=owner,
                stage=stage,
                max_archive_entries=max_archive_entries,
                max_archive_depth=max_archive_depth,
                max_file_size_bytes=max_file_size_bytes,
            )
        )
    return docs


def parse_path(
        path: str | Path,
        *,
        owner: str = OWNER_UNKNOWN,
        stage: str = STAGE_GENERAL,
        max_archive_entries: int = DEFAULT_MAX_ARCHIVE_ENTRIES,
        max_archive_depth: int = DEFAULT_MAX_ARCHIVE_DEPTH,
        max_file_size_bytes: int = DEFAULT_MAX_FILE_SIZE_BYTES,
) -> list[ParsedDocument]:
    path = Path(path)
    data = path.read_bytes()
    mime_type, _ = mimetypes.guess_type(str(path))
    return parse_bytes(
        data,
        filename=path.name,
        mime_type=mime_type,
        owner=owner,
        stage=stage,
        source_name=str(path),
        max_archive_entries=max_archive_entries,
        max_archive_depth=max_archive_depth,
        max_file_size_bytes=max_file_size_bytes,
    )


def parse_bytes(
        data: bytes,
        *,
        filename: str,
        mime_type: Optional[str] = None,
        owner: str = OWNER_UNKNOWN,
        stage: str = STAGE_GENERAL,
        source_name: Optional[str] = None,
        archive_path: Optional[str] = None,
        _archive_depth: int = 0,
        max_archive_entries: int = DEFAULT_MAX_ARCHIVE_ENTRIES,
        max_archive_depth: int = DEFAULT_MAX_ARCHIVE_DEPTH,
        max_file_size_bytes: int = DEFAULT_MAX_FILE_SIZE_BYTES,
) -> list[ParsedDocument]:
    """
    Главная точка разбора.
    Возвращает список документов, чтобы одинаково обрабатывать обычные файлы и архивы.
    """
    filename = filename or "unnamed.bin"
    ext = _detect_extension(filename, mime_type)
    size_bytes = len(data)
    sha256 = _sha256(data)
    mime_type = mime_type or mimetypes.guess_type(filename)[0]

    if size_bytes > max_file_size_bytes:
        return [
            ParsedDocument(
                filename=filename,
                extension=ext,
                text="",
                owner=owner,
                stage=stage,
                parser_name="size_guard",
                mime_type=mime_type,
                source_name=source_name,
                archive_path=archive_path,
                size_bytes=size_bytes,
                sha256=sha256,
                warnings=[f"Файл пропущен: размер превышает лимит {max_file_size_bytes} байт"],
                is_supported=False,
            )
        ]

    if ext == ".zip":
        return _parse_zip(
            data=data,
            filename=filename,
            owner=owner,
            stage=stage,
            source_name=source_name,
            archive_path=archive_path,
            archive_depth=_archive_depth,
            max_archive_entries=max_archive_entries,
            max_archive_depth=max_archive_depth,
            max_file_size_bytes=max_file_size_bytes,
        )

    text, parser_name, warnings, is_supported = _extract_text_by_extension(
        data=data,
        filename=filename,
        mime_type=mime_type,
    )

    doc = ParsedDocument(
        filename=filename,
        extension=ext,
        text=_normalize_text(text),
        owner=owner,
        stage=stage,
        parser_name=parser_name,
        mime_type=mime_type,
        source_name=source_name,
        archive_path=archive_path,
        size_bytes=size_bytes,
        sha256=sha256,
        warnings=warnings,
        is_supported=is_supported,
    )
    return [doc]


def build_context_bundle(
        documents: Iterable[ParsedDocument],
        *,
        max_chars_per_doc: int = DEFAULT_MAX_TEXT_CHARS_PER_DOC,
        max_total_chars: int = DEFAULT_MAX_TOTAL_CONTEXT_CHARS,
        include_empty: bool = False,
        include_warnings: bool = True,
        deduplicate_by_hash: bool = True,
) -> str:
    """
    Собирает единый текстовый контекст для LLM.
    Это пригодится и для согласования темы, и для генерации методички, и для анализа работы.
    """
    seen: set[str] = set()
    parts: list[str] = []
    total = 0

    for doc in documents:
        if deduplicate_by_hash and doc.sha256:
            if doc.sha256 in seen:
                continue
            seen.add(doc.sha256)

        if not include_empty and not (doc.text or "").strip():
            continue

        block = doc.to_context_block(max_chars=max_chars_per_doc)
        if not include_warnings and doc.warnings:
            block = _strip_warning_line(block)

        if total + len(block) > max_total_chars:
            remaining = max_total_chars - total
            if remaining <= 0:
                break
            block = block[:remaining].rstrip() + "\n...[контекст обрезан по общему лимиту]"
            parts.append(block)
            total += len(block)
            break

        parts.append(block)
        total += len(block)

    return "\n\n" + ("\n\n" + ("-" * 80) + "\n\n").join(parts) if parts else ""


def summarize_documents(documents: Iterable[ParsedDocument]) -> dict[str, Any]:
    docs = list(documents)
    by_owner: dict[str, int] = {}
    by_stage: dict[str, int] = {}
    by_ext: dict[str, int] = {}
    warnings_count = 0

    for doc in docs:
        by_owner[doc.owner] = by_owner.get(doc.owner, 0) + 1
        by_stage[doc.stage] = by_stage.get(doc.stage, 0) + 1
        by_ext[doc.extension] = by_ext.get(doc.extension, 0) + 1
        warnings_count += len(doc.warnings)

    return {
        "total_documents": len(docs),
        "total_text_chars": sum(len(doc.text) for doc in docs),
        "total_size_bytes": sum(doc.size_bytes for doc in docs),
        "warnings_count": warnings_count,
        "by_owner": by_owner,
        "by_stage": by_stage,
        "by_extension": by_ext,
    }


def group_documents(
        documents: Iterable[ParsedDocument],
) -> dict[str, dict[str, list[ParsedDocument]]]:
    """
    Возвращает вложенную структуру:
    grouped[owner][stage] -> list[ParsedDocument]
    """
    grouped: dict[str, dict[str, list[ParsedDocument]]] = {}
    for doc in documents:
        grouped.setdefault(doc.owner, {}).setdefault(doc.stage, []).append(doc)
    return grouped


# Legacy-friendly aliases
parse_file = parse_uploaded_file
parse_files = parse_uploaded_files
render_documents_for_llm = build_context_bundle
build_context = build_context_bundle


# Core parser routing


def _extract_text_by_extension(
        *,
        data: bytes,
        filename: str,
        mime_type: Optional[str],
) -> tuple[str, str, list[str], bool]:
    ext = _detect_extension(filename, mime_type)
    warnings: list[str] = []

    try:
        if ext in {".txt", ".md", ".rst", ".log", ".ini", ".cfg", ".yml", ".yaml"}:
            return _parse_text_like(data), "text", warnings, True

        if ext in {".json"}:
            return _parse_json(data), "json", warnings, True

        if ext in {".csv", ".tsv"}:
            return _parse_csv(data, delimiter="\t" if ext == ".tsv" else ","), "csv", warnings, True

        if ext in {".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cpp", ".c", ".h", ".hpp", ".sql"}:
            return _parse_code_like(data), "code", warnings, True

        if ext == ".ipynb":
            return _parse_ipynb(data), "ipynb", warnings, True

        if ext == ".pdf":
            text, local_warnings = _parse_pdf(data)
            warnings.extend(local_warnings)
            return text, "pdf", warnings, True

        if ext == ".docx":
            text, local_warnings = _parse_docx(data)
            warnings.extend(local_warnings)
            return text, "docx", warnings, True

        if ext == ".pptx":
            text, local_warnings = _parse_pptx(data)
            warnings.extend(local_warnings)
            return text, "pptx", warnings, True

        if ext in {".rtf"}:
            return _parse_rtf(data), "rtf", warnings, True

        if ext in {
            ".png", ".jpg", ".jpeg", ".webp", ".bmp", ".gif", ".tif", ".tiff",
            ".xlsx", ".xls", ".ppt", ".doc", ".exe", ".dll", ".so", ".bin",
            ".sqlite", ".db"
        }:
            warnings.append("Файл не был преобразован в текст: неподдерживаемый или бинарный формат")
            return "", "unsupported_binary", warnings, False

        # Фоллбек: пробуем декодировать как обычный текст
        fallback_text = _best_effort_decode(data)
        if fallback_text.strip():
            warnings.append("Формат не распознан точно; применен текстовый Фоллбек")
            return fallback_text, "fallback_text", warnings, True

        warnings.append("Файл не удалось интерпретировать как текст")
        return "", "unknown", warnings, False

    except Exception as exc:
        warnings.append(f"Ошибка разбора: {type(exc).__name__}: {exc}")
        return "", "error", warnings, False


# Specific parsers


def _parse_text_like(data: bytes) -> str:
    return _best_effort_decode(data)


def _parse_code_like(data: bytes) -> str:
    return _best_effort_decode(data)


def _parse_json(data: bytes) -> str:
    raw = _best_effort_decode(data)
    try:
        obj = json.loads(raw)
        return json.dumps(obj, ensure_ascii=False, indent=2)
    except Exception:
        return raw


def _parse_csv(data: bytes, delimiter: str = ",") -> str:
    text = _best_effort_decode(data)
    reader = csv.reader(io.StringIO(text), delimiter=delimiter)
    rows: list[str] = []
    for row in reader:
        rows.append(" | ".join(cell.strip() for cell in row))
    return "\n".join(rows)


def _parse_ipynb(data: bytes) -> str:
    payload = json.loads(_best_effort_decode(data))
    cells = payload.get("cells", [])
    parts: list[str] = []

    for idx, cell in enumerate(cells, start=1):
        cell_type = cell.get("cell_type", "unknown")
        source = "".join(cell.get("source", []))
        source = source.strip()

        if not source:
            continue

        if cell_type == "markdown":
            parts.append(f"[Markdown cell {idx}]\n{source}")
        elif cell_type == "code":
            parts.append(f"[Code cell {idx}]\n{source}")
        else:
            parts.append(f"[{cell_type} cell {idx}]\n{source}")

    return "\n\n".join(parts)


def _parse_pdf(data: bytes) -> tuple[str, list[str]]:
    warnings: list[str] = []
    if PdfReader is None:
        return "", ["Не установлен пакет для разбора PDF (pypdf или PyPDF2)"]

    reader = PdfReader(io.BytesIO(data))
    parts: list[str] = []

    total_pages = len(reader.pages)
    for idx, page in enumerate(reader.pages, start=1):
        try:
            page_text = page.extract_text() or ""
        except Exception as exc:
            page_text = ""
            warnings.append(f"Не удалось извлечь текст со страницы {idx}: {type(exc).__name__}")

        page_text = _normalize_text(page_text)
        if page_text.strip():
            parts.append(f"[Страница {idx} из {total_pages}]\n{page_text}")

    if not parts:
        warnings.append("PDF разобран, но извлеченный текст пустой; возможно, документ состоит из изображений")
    return "\n\n".join(parts), warnings


def _parse_docx(data: bytes) -> tuple[str, list[str]]:
    warnings: list[str] = []
    if DocxDocument is None:
        return "", ["Не установлен пакет python-docx"]

    document = DocxDocument(io.BytesIO(data))
    parts: list[str] = []

    # Paragraphs
    for p in document.paragraphs:
        text = (p.text or "").strip()
        if text:
            parts.append(text)

    # Tables
    for table_idx, table in enumerate(document.tables, start=1):
        table_rows: list[str] = []
        for row in table.rows:
            cells = [(cell.text or "").strip() for cell in row.cells]
            if any(cells):
                table_rows.append(" | ".join(cells))
        if table_rows:
            parts.append(f"[Таблица {table_idx}]\n" + "\n".join(table_rows))

    if not parts:
        warnings.append("DOCX разобран, но текст не найден")
    return "\n\n".join(parts), warnings


def _parse_pptx(data: bytes) -> tuple[str, list[str]]:
    warnings: list[str] = []
    if Presentation is None:
        return "", ["Не установлен пакет python-pptx"]

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_parts: list[str] = []
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text and str(text).strip():
                slide_parts.append(str(text).strip())
        if slide_parts:
            parts.append(f"[Слайд {slide_idx}]\n" + "\n".join(slide_parts))

    if not parts:
        warnings.append("PPTX разобран, но текст на слайдах не найден")
    return "\n\n".join(parts), warnings


def _parse_rtf(data: bytes) -> str:
    """
    Простой и безопасный фоллбек для RTF без внешних зависимостей.
    Не идеален, но для текстовых методичек и заметок обычно достаточен.
    """
    text = _best_effort_decode(data)

    # Удаляем управляющие слова RTF грубым способом
    text = re.sub(r"\\par[d]?", "\n", text)
    text = re.sub(r"\\'[0-9a-fA-F]{2}", "", text)
    text = re.sub(r"\\[a-zA-Z]+\d* ?", "", text)
    text = text.replace("{", "").replace("}", "")
    return _normalize_text(text)


def _parse_zip(
        *,
        data: bytes,
        filename: str,
        owner: str,
        stage: str,
        source_name: Optional[str],
        archive_path: Optional[str],
        archive_depth: int,
        max_archive_entries: int,
        max_archive_depth: int,
        max_file_size_bytes: int,
) -> list[ParsedDocument]:
    docs: list[ParsedDocument] = []

    if archive_depth >= max_archive_depth:
        return [
            ParsedDocument(
                filename=filename,
                extension=".zip",
                text="",
                owner=owner,
                stage=stage,
                parser_name="zip_guard",
                mime_type="application/zip",
                source_name=source_name,
                archive_path=archive_path,
                size_bytes=len(data),
                sha256=_sha256(data),
                warnings=[f"Архив пропущен: превышена глубина вложенности ({max_archive_depth})"],
                is_supported=False,
            )
        ]

    current_archive_path = f"{archive_path}!{filename}" if archive_path else filename

    try:
        with zipfile.ZipFile(io.BytesIO(data)) as zf:
            names = [n for n in zf.namelist() if not n.endswith("/")]
            names = [n for n in names if not _is_probably_system_file(n)]

            if len(names) > max_archive_entries:
                names = names[:max_archive_entries]
                archive_warning = f"Архив обрезан: обработаны только первые {max_archive_entries} файлов"
            else:
                archive_warning = ""

            for member_name in names:
                try:
                    member_data = zf.read(member_name)
                except Exception as exc:
                    docs.append(
                        ParsedDocument(
                            filename=Path(member_name).name or member_name,
                            extension=_detect_extension(member_name, None),
                            text="",
                            owner=owner,
                            stage=stage,
                            parser_name="zip_member_error",
                            mime_type=mimetypes.guess_type(member_name)[0],
                            source_name=source_name,
                            archive_path=current_archive_path,
                            size_bytes=0,
                            sha256="",
                            warnings=[f"Не удалось прочитать файл из архива: {type(exc).__name__}: {exc}"],
                            is_supported=False,
                        )
                    )
                    continue

                child_docs = parse_bytes(
                    member_data,
                    filename=Path(member_name).name or member_name,
                    mime_type=mimetypes.guess_type(member_name)[0],
                    owner=owner,
                    stage=stage,
                    source_name=source_name,
                    archive_path=current_archive_path,
                    _archive_depth=archive_depth + 1,
                    max_archive_entries=max_archive_entries,
                    max_archive_depth=max_archive_depth,
                    max_file_size_bytes=max_file_size_bytes,
                )

                # Сохраняем относительный путь внутри архива в предупреждениях
                for child in child_docs:
                    rel_path = member_name.replace("\\", "/")
                    child.warnings = [f"Путь в архиве: {rel_path}"] + child.warnings
                    if archive_warning and archive_warning not in child.warnings:
                        child.warnings.append(archive_warning)

                docs.extend(child_docs)

    except zipfile.BadZipFile:
        docs.append(
            ParsedDocument(
                filename=filename,
                extension=".zip",
                text="",
                owner=owner,
                stage=stage,
                parser_name="bad_zip",
                mime_type="application/zip",
                source_name=source_name,
                archive_path=archive_path,
                size_bytes=len(data),
                sha256=_sha256(data),
                warnings=["Файл имеет расширение .zip, но не является корректным архивом"],
                is_supported=False,
            )
        )

    return docs


# Helpers


def _coerce_input(obj: Any) -> tuple[bytes, str, Optional[str], Optional[str]]:
    """
    Возвращает:
    - bytes
    - filename
    - mime_type
    - source_name
    """
    # pathlib / str path
    if isinstance(obj, (str, Path)) and os.path.exists(str(obj)):
        path = Path(obj)
        return (
            path.read_bytes(),
            path.name,
            mimetypes.guess_type(str(path))[0],
            str(path),
        )

    # raw bytes
    if isinstance(obj, (bytes, bytearray)):
        return bytes(obj), "uploaded.bin", None, None

    # Streamlit UploadedFile часто имеет .name, .type, .getvalue()
    if hasattr(obj, "getvalue"):
        data = obj.getvalue()
        filename = getattr(obj, "name", None) or "uploaded.bin"
        mime_type = getattr(obj, "type", None)
        source_name = getattr(obj, "name", None)
        return data, filename, mime_type, source_name

    # Generic file-like
    if hasattr(obj, "read"):
        pos = None
        try:
            if hasattr(obj, "tell"):
                pos = obj.tell()
        except Exception:
            pos = None

        data = obj.read()
        if isinstance(data, str):
            data = data.encode("utf-8", errors="ignore")

        if pos is not None and hasattr(obj, "seek"):
            try:
                obj.seek(pos)
            except Exception:
                pass

        filename = getattr(obj, "name", None) or "uploaded.bin"
        mime_type = mimetypes.guess_type(str(filename))[0]
        source_name = getattr(obj, "name", None)
        return bytes(data), Path(str(filename)).name, mime_type, source_name

    raise TypeError(f"Неподдерживаемый тип входа для разбора: {type(obj).__name__}")


def _detect_extension(filename: str, mime_type: Optional[str]) -> str:
    ext = Path(filename).suffix.lower()
    if ext:
        return ext

    if mime_type:
        guessed = mimetypes.guess_extension(mime_type)
        if guessed:
            return guessed.lower()

    return ""


def _best_effort_decode(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "cp1251", "latin-1"):
        try:
            return data.decode(enc)
        except Exception:
            continue
    return data.decode("utf-8", errors="ignore")


def _normalize_text(text: str) -> str:
    text = text.replace("\x00", " ")
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+\n", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


def _strip_warning_line(block: str) -> str:
    lines = block.splitlines()
    filtered = [line for line in lines if not line.startswith("Предупреждения:")]
    return "\n".join(filtered)


def _sha256(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def _is_probably_system_file(path_in_archive: str) -> bool:
    normalized = path_in_archive.replace("\\", "/")
    basename = Path(normalized).name

    if basename.startswith("."):
        return True

    ignored_fragments = [
        "__macosx/",
        ".ds_store",
        "thumbs.db",
        ".git/",
        ".idea/",
        ".vscode/",
    ]
    lower_path = normalized.lower()
    return any(fragment in lower_path for fragment in ignored_fragments)


# Demo / manual smoke test


if __name__ == "__main__":
    import sys

    if len(sys.argv) < 2:
        print("Usage: python parsers.py <file-or-archive> [owner] [stage]")
        raise SystemExit(1)

    input_path = sys.argv[1]
    owner = sys.argv[2] if len(sys.argv) > 2 else OWNER_UNKNOWN
    stage = sys.argv[3] if len(sys.argv) > 3 else STAGE_GENERAL

    docs = parse_path(input_path, owner=owner, stage=stage)
    print(json.dumps(summarize_documents(docs), ensure_ascii=False, indent=2))
    print()
    print(build_context_bundle(docs, max_chars_per_doc=2500, max_total_chars=10000))
